import sys
from banco_dados.conexao import conecta
from banco_dados.controle_erros import grava_erro_banco
from dados_email import email_user, password
import os
import traceback
import inspect

from openpyxl import Workbook
from openpyxl.styles import Alignment, Font, PatternFill, Border, Side
from openpyxl.utils import get_column_letter
from datetime import date, datetime


import re
import fitz
from openpyxl.drawing.image import Image as XLImage
from io import BytesIO

import smtplib
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from email.mime.base import MIMEBase
from email.header import Header
from email import encoders


class ClassificarOps:
    def __init__(self):
        nome_arquivo_com_caminho = inspect.getframeinfo(inspect.currentframe()).filename
        self.nome_arquivo = os.path.basename(nome_arquivo_com_caminho)
        self.diretorio_script = os.path.dirname(nome_arquivo_com_caminho)
        nome_base = os.path.splitext(self.nome_arquivo)[0]
        self.arquivo_log = os.path.join(self.diretorio_script, f"{nome_base}_erros.txt")

        self.manipula_comeco()

    def trata_excecao(self, nome_funcao, mensagem, arquivo, excecao):
        try:
            tb = traceback.extract_tb(excecao)
            num_linha_erro = tb[-1][1]

            traceback.print_exc()
            print(f'Houve um problema no arquivo: {arquivo} na função: "{nome_funcao}"\n{mensagem} {num_linha_erro}')

            grava_erro_banco(nome_funcao, mensagem, arquivo, num_linha_erro)

            # 'Log' em arquivo local apenas se houver erro
            with open(self.arquivo_log, "a", encoding="utf-8") as f:
                f.write(f"Erro na função {nome_funcao} do arquivo {arquivo}: {mensagem} (linha {num_linha_erro})\n")

        except Exception as e:
            nome_funcao_trat = inspect.currentframe().f_code.co_name
            exc_traceback = sys.exc_info()[2]
            tb = traceback.extract_tb(exc_traceback)
            num_linha_erro = tb[-1][1]
            print(f'Houve um problema no arquivo: {self.nome_arquivo} na função: "{nome_funcao_trat}"\n'
                  f'{e} {num_linha_erro}')
            grava_erro_banco(nome_funcao_trat, e, self.nome_arquivo, num_linha_erro)

            with open(self.arquivo_log, "a", encoding="utf-8") as f:
                f.write(
                    f"Erro na função {nome_funcao_trat} do arquivo {self.nome_arquivo}: {e} (linha {num_linha_erro})\n")

    def dados_email(self):
        try:
            to = ['<maquinas@unisold.com.br>', '<ahcmaquinas@gmail.com>']

            current_time = (datetime.now())
            horario = current_time.strftime('%H')
            hora_int = int(horario)
            saudacao = ""
            if 4 < hora_int < 13:
                saudacao = "Bom Dia!"
            elif 12 < hora_int < 19:
                saudacao = "Boa Tarde!"
            elif hora_int > 18:
                saudacao = "Boa Noite!"
            elif hora_int < 5:
                saudacao = "Boa Noite!"

            msg_final = f"Att,\n" \
                        f"Suzuki Máquinas Ltda\n" \
                        f"Fone (51) 3561.2583/(51) 3170.0965\n\n" \
                        f"Mensagem enviada automaticamente, por favor não responda.\n\n" \
                        f"Se houver algum problema com o recebimento de emails ou conflitos com o arquivo excel, " \
                        f"favor entrar em contato pelo email maquinas@unisold.com.br.\n\n"

            return saudacao, msg_final, to

        except Exception as e:
            nome_funcao = inspect.currentframe().f_code.co_name
            exc_traceback = sys.exc_info()[2]
            self.trata_excecao(nome_funcao, str(e), self.nome_arquivo, exc_traceback)

    def envia_email_excel(self, arquivo):
        try:
            saudacao, msg_final, to = self.dados_email()

            to = ['<maquinas@unisold.com.br>']

            subject = f'Classificação das Ordens de Produção'

            msg = MIMEMultipart()
            msg['From'] = email_user
            msg['Subject'] = subject

            body = f"{saudacao}\n\n" \
                   f"Segue em anexo lista de localização e classificação das Ordens de Produção.\n\n" \
                   f"{msg_final}"

            msg.attach(MIMEText(body, 'plain'))


            attachment = open(arquivo, 'rb')

            part = MIMEBase('application', "octet-stream")
            part.set_payload(attachment.read())
            encoders.encode_base64(part)
            part.add_header('Content-Disposition', 'attachment', filename=Header(arquivo, 'utf-8').encode())
            msg.attach(part)

            text = msg.as_string()
            server = smtplib.SMTP('smtp.gmail.com', 587)
            server.starttls()
            server.login(email_user, password)

            server.sendmail(email_user, to, text)
            attachment.close()

            server.quit()

            print(f'Email enviada com sucesso!')

        except Exception as e:
            nome_funcao = inspect.currentframe().f_code.co_name
            exc_traceback = sys.exc_info()[2]
            self.trata_excecao(nome_funcao, str(e), self.nome_arquivo, exc_traceback)

    def excluir_arquivo(self, caminho_arquivo):
        try:
            if os.path.exists(caminho_arquivo):
                os.remove(caminho_arquivo)
            else:
                print("O arquivo não existe no caminho especificado.")

        except Exception as e:
            nome_funcao = inspect.currentframe().f_code.co_name
            exc_traceback = sys.exc_info()[2]
            self.trata_excecao(nome_funcao, str(e), self.nome_arquivo, exc_traceback)

    def manipula_comeco(self):
        try:
            previsao = datetime.now()

            cursor = conecta.cursor()
            cursor.execute(f"SELECT prod.id, op.numero, op.codigo, op.id_estrutura, prod.descricao, "
                           f"COALESCE(prod.obs, ''), "
                           f"prod.unidade, COALESCE(prod.tipomaterial, ''), op.quantidade, "
                           f"COALESCE(ser.descricao, '') "
                           f"FROM ordemservico as op "
                           f"INNER JOIN produto as prod ON op.produto = prod.id "
                           f"LEFT JOIN SERVICO_INTERNO as ser ON ser.id = prod.id_servico_interno "
                           f"where op.status = 'A' order by op.numero;")
            ops_abertas = cursor.fetchall()

            ops_por_numero = {}

            if ops_abertas:
                for i in ops_abertas:
                    id_produto, num_op, cod, id_estrutura, descr, ref, um, tipo, qtde, servico_in = i

                    cursor = conecta.cursor()
                    cursor.execute(f"SELECT id, id_pedidointerno, ID_PRODUTO_PI, tipo, numero, id_produto "
                                   f"FROM VINCULO_PRODUTO_PI "
                                   f"where numero = '{num_op}' and tipo = 'OP' "
                                   f"and id_produto = {id_produto};")
                    consulta_vinculos = cursor.fetchall()
                    if consulta_vinculos:
                        for ii in consulta_vinculos:
                            id_vinculo, id_pedido, id_produto_pedido, tipo, numero, id_produto = ii

                            cursor = conecta.cursor()
                            cursor.execute(f"SELECT ped.emissao, prod.codigo, "
                                           f"prod.descricao, "
                                           f"COALESCE(prod.obs, '') as obs, "
                                           f"prod.unidade, prodint.data_previsao "
                                           f"FROM PRODUTOPEDIDOINTERNO as prodint "
                                           f"INNER JOIN produto as prod ON prodint.id_produto = prod.id "
                                           f"INNER JOIN pedidointerno as ped ON prodint.id_pedidointerno = ped.id "
                                           f"INNER JOIN clientes as cli ON ped.id_cliente = cli.id "
                                           f"where prodint.status = 'A' "
                                           f"and prodint.id_produto = {id_produto_pedido} "
                                           f"and prodint.id_pedidointerno = {id_pedido} "
                                           f"ORDER BY prodint.data_previsao ASC;")
                            dados_interno = cursor.fetchall()
                            if dados_interno:
                                for iii in dados_interno:
                                    emissao, codi, descri, refi, umi, previsao = iii

                                    dados = (num_op, id_pedido, cod, descr, ref, um, qtde, previsao, servico_in)

                                    # Se ainda não existe essa OP no dicionário
                                    if num_op not in ops_por_numero:
                                        ops_por_numero[num_op] = dados
                                    else:
                                        # Compara a data de previsão (posição 5 da tupla)
                                        previsao_atual = ops_por_numero[num_op][7]
                                        if previsao < previsao_atual:
                                            ops_por_numero[num_op] = dados

                    else:
                        dados = (num_op, "sem vinculo", cod, descr, ref, um, qtde, previsao, servico_in)

                        # Se ainda não existe essa OP no dicionário
                        if num_op not in ops_por_numero:
                            ops_por_numero[num_op] = dados
                        else:
                            # Compara a data de previsão (posição 5 da tupla)
                            previsao_atual = ops_por_numero[num_op][7]
                            if previsao < previsao_atual:
                                ops_por_numero[num_op] = dados

                lista = list(ops_por_numero.values())

                if lista:
                    self.manipula_dados(lista)

        except Exception as e:
            nome_funcao = inspect.currentframe().f_code.co_name
            exc_traceback = sys.exc_info()[2]
            self.trata_excecao(nome_funcao, str(e), self.nome_arquivo, exc_traceback)

    def manipula_dados(self, lista):
        try:
            op_ab_editado = []

            for i in lista:
                num_op, pi, cod, descr, ref, um, qtde, previsao, servico = i

                cursor = conecta.cursor()
                cursor.execute(f"select ordser.id, ordser.datainicial, ordser.dataprevisao, ordser.numero, prod.codigo, "
                               f"prod.descricao, "
                               f"COALESCE(prod.obs, '') as obs, prod.unidade, "
                               f"ordser.quantidade, ordser.id_estrutura, ordser.LOCAL_OP, prod.TIPOMATERIAL "
                               f"from ordemservico as ordser "
                               f"INNER JOIN produto prod ON ordser.produto = prod.id "
                               f"where ordser.status = 'A' and ordser.numero = '{num_op}';")
                op_abertas = cursor.fetchall()
                if op_abertas:
                    id_op, emissao, previsao_aaa, op, cod, descr, ref, um, qtde, id_estrut, local_op, tipo = op_abertas[0]

                    if id_estrut:
                        total_estrut = 0
                        total_consumo = 0

                        cursor = conecta.cursor()
                        cursor.execute(f"SELECT estprod.id, "
                                       f"((SELECT quantidade FROM ordemservico where numero = {op}) * "
                                       f"(estprod.quantidade)) AS Qtde "
                                       f"FROM estrutura_produto as estprod "
                                       f"INNER JOIN produto prod ON estprod.id_prod_filho = prod.id "
                                       f"where estprod.id_estrutura = {id_estrut};")
                        itens_estrutura = cursor.fetchall()

                        for dads in itens_estrutura:
                            ides, quantidade = dads
                            total_estrut += 1

                            cursor = conecta.cursor()
                            cursor.execute(f"SELECT max(prodser.ID_ESTRUT_PROD), "
                                           f"sum(prodser.QTDE_ESTRUT_PROD) as total "
                                           f"FROM estrutura_produto as estprod "
                                           f"INNER JOIN produto prod ON estprod.id_prod_filho = prod.id "
                                           f"INNER JOIN produtoos as prodser ON estprod.id = prodser.ID_ESTRUT_PROD "
                                           f"where prodser.numero = {op} and estprod.id = {ides} "
                                           f"group by prodser.ID_ESTRUT_PROD;")
                            itens_consumo = cursor.fetchall()
                            for duds in itens_consumo:
                                id_mats, qtde_mats = duds
                                if ides == id_mats and quantidade == qtde_mats:
                                    total_consumo += 1

                        msg = f"{total_estrut}/{total_consumo}"

                        esta_na_pasta = self.verifica_arquivos_aguardando(num_op)

                        if esta_na_pasta:
                            cursor = conecta.cursor()
                            cursor.execute(f"UPDATE ordemservico SET LOCAL_OP = 'ALMOX' "
                                           f"WHERE id = {id_op};")

                            conecta.commit()

                            dados = (num_op, pi, cod, descr, ref, um, qtde, previsao, servico, msg, "ALMOX")
                            op_ab_editado.append(dados)
                        elif local_op == "PROJETO":
                            dados = (num_op, pi, cod, descr, ref, um, qtde, previsao, servico, msg, local_op)
                            op_ab_editado.append(dados)
                        elif local_op == "ALMOX":
                            dados = (num_op, pi, cod, descr, ref, um, qtde, previsao, servico, msg, local_op)
                            op_ab_editado.append(dados)
                        elif not servico:
                            dados = (num_op, pi, cod, descr, ref, um, qtde, previsao, servico, msg,
                                     "SEM SERVIÇO DEFINIDO")
                            op_ab_editado.append(dados)
                        elif servico != "MONTAGEM" and servico != "SOLDA" and servico != "ELETRICO":
                            if total_consumo == 0:
                                dados = (num_op, pi, cod, descr, ref, um, qtde, previsao, servico, msg,
                                         "LOC. CORTE")
                                op_ab_editado.append(dados)
                            else:
                                dados = (num_op, pi, cod, descr, ref, um, qtde, previsao, servico, msg,
                                         "LOC. USINAGEM")
                                op_ab_editado.append(dados)
                        elif servico == "SOLDA":
                            dados = (num_op, pi, cod, descr, ref, um, qtde, previsao, servico, msg, "LOC. SOLDA")
                            op_ab_editado.append(dados)
                        elif servico == "ELETRICO":
                            dados = (num_op, pi, cod, descr, ref, um, qtde, previsao, servico, msg, "LOC. ELETRICA")
                            op_ab_editado.append(dados)
                        elif servico == "MONTAGEM":
                            dados = (num_op, pi, cod, descr, ref, um, qtde, previsao, servico, msg,
                                     "LOC. MONTAGEM")
                            op_ab_editado.append(dados)
                        else:
                            dados = (num_op, pi, cod, descr, ref, um, qtde, previsao, servico, msg,
                                     "NÃO SEI")
                            op_ab_editado.append(dados)

            if op_ab_editado:
                arquivo = "ops_abertas.xlsx"

                self.gerar_excel(arquivo, op_ab_editado)
                self.envia_email_excel(arquivo)
                self.excluir_arquivo(arquivo)
                #self.gerar_power_point(op_ab_editado)

        except Exception as e:
            nome_funcao = inspect.currentframe().f_code.co_name
            exc_traceback = sys.exc_info()[2]
            self.trata_excecao(nome_funcao, str(e), self.nome_arquivo, exc_traceback)

    def gerar_excel(self, arquivo, lista):
        try:
            wb = Workbook()
            wb.remove(wb.active)

            cabecalhos = [
                "Nº OP", "Nº PI", "Código", "Descrição", "Referência",
                "UM", "Qtde", "Previsão", "Serviço Int.", "Consumo", "LOCAL_OP",
                "Imagem"
            ]

            alinhamento_central = Alignment(horizontal="center", vertical="center")

            borda = Border(
                left=Side(style="thin"),
                right=Side(style="thin"),
                top=Side(style="thin"),
                bottom=Side(style="thin")
            )

            fonte_negrito = Font(bold=True)

            fundo_cinza = PatternFill(
                start_color="DDDDDD",
                end_color="DDDDDD",
                fill_type="solid"
            )

            abas = {}

            # 🔹 Separar dados por aba
            for linha in lista:
                num_op, pi, cod, descr, ref, um, qtde, previsao, servico, consumo, local_op = linha

                if local_op and local_op.upper() == "PROJETO":
                    nome_aba = "PROJETO"
                elif local_op and local_op.upper() == "ALMOX":
                    nome_aba = "ALMOX"
                else:
                    nome_aba = local_op if local_op else "SEM SERVIÇO"

                nome_aba = str(nome_aba)[:31]

                if nome_aba not in abas:
                    abas[nome_aba] = []

                abas[nome_aba].append(linha)

            # 🔹 Criar abas
            for nome_aba, dados_aba in abas.items():

                dados_aba.sort(
                    key=lambda x: x[7] if isinstance(x[7], (date, datetime)) else datetime.max
                )

                ws = wb.create_sheet(title=nome_aba)
                ws.append(cabecalhos)

                # Cabeçalho estilizado
                for cell in ws[1]:
                    cell.font = fonte_negrito
                    cell.fill = fundo_cinza
                    cell.alignment = alinhamento_central
                    cell.border = borda

                # 🔹 Inserir dados
                for linha in dados_aba:
                    num_op, pi, cod, descr, ref, um, qtde, previsao, servico, consumo, local_op = linha

                    ws.append([
                        num_op,
                        pi,
                        cod,
                        descr,
                        ref,
                        um,
                        qtde,
                        previsao,
                        servico,
                        consumo,
                        local_op,
                        ""
                    ])

                    linha_excel = ws.max_row

                    # 🔹 Gerar imagem do PDF em memória
                    s = re.sub(r"[^\d.]", "", str(ref))
                    s = re.sub(r"\.+$", "", s)

                    caminho_pdf = rf"\\Publico\C\OP\Projetos\{s}.pdf"

                    if os.path.exists(caminho_pdf):

                        try:
                            from PIL import Image
                            doc = fitz.open(caminho_pdf)
                            page = doc.load_page(0)
                            pix = page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5))

                            # Converte pixmap para imagem PIL
                            img_pil = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)

                            # Salva como JPEG com compressão
                            img_bytes = BytesIO()
                            img_pil.save(img_bytes, format="JPEG", quality=40, optimize=True)
                            img_bytes.seek(0)

                            doc.close()

                            img = XLImage(img_bytes)
                            img.width = 120
                            img.height = 90

                            coluna_imagem = ws.max_column
                            celula = f"{get_column_letter(coluna_imagem)}{linha_excel}"

                            ws.add_image(img, celula)

                            ws.row_dimensions[linha_excel].height = 70

                        except:
                            pass  # Se der erro no PDF, ignora e segue

                # 🔹 Formatação dados (sem coluna imagem)
                for row in ws.iter_rows(min_row=2, max_row=ws.max_row, max_col=11):
                    for cell in row:
                        cell.alignment = alinhamento_central
                        cell.border = borda

                        if isinstance(cell.value, (date, datetime)):
                            cell.number_format = "DD/MM/YYYY"

                # 🔹 Ajuste largura automática (exceto imagem)
                for col in ws.iter_cols(min_col=1, max_col=11):
                    max_length = 0
                    col_letter = get_column_letter(col[0].column)

                    for cell in col:
                        if cell.value:
                            max_length = max(max_length, len(str(cell.value)))

                    ws.column_dimensions[col_letter].width = max_length + 2

                # 🔹 Largura fixa da coluna imagem
                ws.column_dimensions[get_column_letter(ws.max_column)].width = 20

            wb.save(arquivo)

            print("Excel Gerado com imagens!")

        except Exception as e:
            nome_funcao = inspect.currentframe().f_code.co_name
            exc_traceback = sys.exc_info()[2]
            self.trata_excecao(nome_funcao, str(e), self.nome_arquivo, exc_traceback)

    def gerar_power_point(self, lista):
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
            from pptx.enum.text import MSO_AUTO_SIZE
            from datetime import date, datetime
            import os
            import re
            import fitz
            from io import BytesIO

            desktop = os.path.join(os.path.expanduser("~"), "Desktop")
            arquivo = os.path.join(desktop, "painel_producao_tv.pptx")

            prs = Presentation()
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)

            abas = {}

            # 🔹 Separar por setor (igual seu Excel)
            for linha in lista:
                num_op, pi, cod, descr, ref, um, qtde, previsao, servico, consumo, local_op = linha

                if local_op and local_op.upper() == "PROJETO":
                    nome_aba = "PROJETO"
                elif local_op and local_op.upper() == "ALMOX":
                    nome_aba = "ALMOX"
                else:
                    nome_aba = local_op if local_op else "SEM SERVIÇO"

                nome_aba = str(nome_aba)[:31]

                if nome_aba not in abas:
                    abas[nome_aba] = []

                abas[nome_aba].append(linha)

            # 🔹 Criar slide por setor
            for nome_aba, dados_aba in abas.items():

                # ordenar por previsão
                dados_aba.sort(
                    key=lambda x: x[7] if isinstance(x[7], (date, datetime)) else datetime.max
                )

                # pegar apenas as 4 primeiras prioridades
                dados_aba = dados_aba[:4]

                slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(slide_layout)

                # Título setor
                title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.5))
                tf_title = title.text_frame
                tf_title.text = f"SETOR: {nome_aba}"
                tf_title.paragraphs[0].font.size = Pt(32)
                tf_title.paragraphs[0].font.bold = True

                card_width = Inches(6.1)
                card_height = Inches(3)

                positions = [
                    (Inches(0.5), Inches(1)),
                    (Inches(6.7), Inches(1)),
                    (Inches(0.5), Inches(4.2)),
                    (Inches(6.7), Inches(4.2)),
                ]

                for pos, linha in zip(positions, dados_aba):

                    num_op, pi, cod, descr, ref, um, qtde, previsao, servico, consumo, local_op = linha

                    left, top = pos
                    width = card_width
                    height = card_height

                    # cartão base
                    card = slide.shapes.add_shape(
                        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                        left, top, width, height
                    )
                    card.line.width = Pt(2)

                    # 🔹 inserir imagem do PDF
                    s = re.sub(r"[^\d.]", "", str(ref))
                    s = re.sub(r"\.+$", "", s)
                    caminho_pdf = rf"\\Publico\C\OP\Projetos\{s}.pdf"

                    # 🔹 proporção baseada no tamanho do card
                    altura_imagem = card_height * 0.8
                    altura_texto = card_height * 0.15

                    if os.path.exists(caminho_pdf):
                        try:
                            doc = fitz.open(caminho_pdf)
                            page = doc.load_page(0)
                            pix = page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5))
                            img_bytes = BytesIO(pix.tobytes("png"))
                            doc.close()

                            # 🔹 IMAGEM GRANDE
                            slide.shapes.add_picture(
                                img_bytes,
                                left + Inches(0.15),
                                top + Inches(0.15),
                                width=card_width - Inches(0.3),
                                height=altura_imagem
                            )
                        except:
                            pass

                    # 🔹 TEXTO BEM EMBAIXO
                    info = slide.shapes.add_textbox(
                        left + Inches(0.15),
                        top + Inches(0.15) + altura_imagem + Inches(0.05),
                        card_width - Inches(0.3),
                        altura_texto
                    )

                    tf = info.text_frame
                    tf.clear()

                    texto_prev = previsao.strftime("%d/%m/%Y") if isinstance(previsao, (date, datetime)) else ""

                    p = tf.paragraphs[0]
                    p.text = f"OP: {num_op} | PREV: {texto_prev} | QTDE: {qtde}"
                    p.font.size = Pt(18)
                    p.font.bold = True
                    p.alignment = PP_ALIGN.CENTER

            prs.save(arquivo)

            print("Painel TV gerado com sucesso!")

        except Exception as e:
            nome_funcao = inspect.currentframe().f_code.co_name
            exc_traceback = sys.exc_info()[2]
            self.trata_excecao(nome_funcao, str(e), self.nome_arquivo, exc_traceback)

    def verifica_arquivos_aguardando(self, num_op_bc):
        try:
            esta_na_pasta = False

            arquivos_pdf = self.procura_arquivos_aguarda_material()

            numeros_op_processados = []
            duplicados = []

            for arq_original in arquivos_pdf:
                dadinhos = arq_original[35:]

                inicio = dadinhos.find("OP ")
                dadinhos1 = dadinhos[inicio + 3:]

                ini = dadinhos1.find(" - ")
                num_op = dadinhos1[:ini]

                if num_op in numeros_op_processados:
                    duplicados.append(num_op)
                else:
                    numeros_op_processados.append(num_op)

            if not duplicados:
                for arq_original in arquivos_pdf:
                    arq_original = arq_original[35:]

                    inicio = arq_original.find("OP ")
                    dadinhos1 = arq_original[inicio + 3:]

                    ini = dadinhos1.find(" - ")
                    num_op = dadinhos1[:ini]
                    if str(num_op_bc) == num_op:
                        esta_na_pasta = True
                        break

            return esta_na_pasta

        except Exception as e:
            nome_funcao = inspect.currentframe().f_code.co_name
            exc_traceback = sys.exc_info()[2]
            self.trata_excecao(nome_funcao, str(e), self.nome_arquivo, exc_traceback)

    def procura_arquivos_aguarda_material(self):
        try:
            caminho = r'\\publico\C\OP\Aguardando Material/'
            extensao = ".pdf"

            arquivos_pdfs = []

            for arquivo in os.listdir(caminho):
                if arquivo.endswith(extensao):
                    caminho_arquivo = os.path.join(caminho, arquivo)
                    arquivos_pdfs.append(caminho_arquivo)

            return arquivos_pdfs

        except Exception as e:
            nome_funcao = inspect.currentframe().f_code.co_name
            exc_traceback = sys.exc_info()[2]
            self.trata_excecao(nome_funcao, str(e), self.nome_arquivo, exc_traceback)


chama_classe = ClassificarOps()